import streamlit as st
import docx2txt
import re
import pdfplumber
import pptx
from translation import translate_text


def extract_text_from_txt(file_name):
    with open(file_name) as file:
        text = file.read()

    return text


def extract_text_from_pdf(file):
    with pdfplumber.open(file) as pdf:
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"  # Add a newline after each page

    return text


def extract_text_from_docx(file):
    text = docx2txt.process(file)
    return text


def extract_text_from_ppt(file_path):
    text = ""
    ppt = pptx.Presentation(file_path)

    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text += run.text + "\n"

    return text


def limit_consecutive_newlines(text):
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text


def split_text_into_lines(text, limit):
    lines = []
    current_line = ""

    for char in text:
        current_line += char

        if len(current_line) >= limit and char == " ":
            lines.append(current_line.rstrip())
            current_line = ""
        elif char == "\n":
            lines.append(current_line.rstrip())
            current_line = ""

    if current_line:
        lines.append(current_line.rstrip())

    return "\n".join(lines)


def main():
    st.set_page_config(layout="wide")
    st.title("Language Bridge")

    uploaded_file = st.file_uploader("File", type=["txt", "pdf", "docx", "pptx"], label_visibility="hidden")
    language_dict = {"Assamese": "as", "Hindi": "hi", "Marathi": "mr", "Tamil": "ta", "Telugu": "te"}
    selected_language = st.selectbox("Select target language", [key for key in language_dict])

    if st.button("Translate"):
        if uploaded_file is not None:
            file_type = uploaded_file.name.split(".")[1].lower()

            if selected_language == "Assamese":
                limit = 174
            elif selected_language in ["Hindi", "Marathi"]:
                limit = 184
            elif selected_language == "Tamil":
                limit = 104
            else:
                limit = 154

            if file_type == "txt":
                text = extract_text_from_txt(uploaded_file.name)
                translated_text = limit_consecutive_newlines(translate_text(text, language_dict[selected_language]))
                st.code(split_text_into_lines(translated_text, limit), language="markdown")
                st.download_button(label="Download TXT", data=translated_text.encode("utf-8"),
                                   file_name=f"{uploaded_file.name.split('.')[0]}_{selected_language.lower()}.txt")
            elif file_type == "pdf":
                text = extract_text_from_pdf(uploaded_file)
                translated_text = translate_text(text, language_dict[selected_language])
                st.code(split_text_into_lines(translated_text, limit), language="markdown")
                st.download_button("Download Translated Text", data=translated_text.encode("utf-8"),
                                   file_name=f"{uploaded_file.name.split('.')[0]}_{selected_language.lower()}.txt")
            elif file_type == "docx":
                text = extract_text_from_docx(uploaded_file)
                translated_text = limit_consecutive_newlines(translate_text(text, language_dict[selected_language]))
                st.code(split_text_into_lines(translated_text, limit), language="markdown")
                st.download_button("Download Translated Text", data=translated_text.encode("utf-8"),
                                   file_name=f"{uploaded_file.name.split('.')[0]}_{selected_language.lower()}.txt")
            else:
                text = extract_text_from_ppt(uploaded_file)
                translated_text = translate_text(text, language_dict[selected_language])
                st.code(split_text_into_lines(translated_text, limit), language="markdown")
                st.download_button("Download Translated Text", data=translated_text.encode("utf-8"),
                                   file_name=f"{uploaded_file.name.split('.')[0]}_{selected_language.lower()}.txt")

        else:
            st.warning("Please upload a file")


if __name__ == "__main__":
    main()
